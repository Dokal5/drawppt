from __future__ import annotations

import uuid
from pathlib import Path
from typing import Literal

from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore


app = FastAPI(title="DrawPPT API", version="0.1.0")

EXPORT_DIR = Path("exports")
EXPORT_DIR.mkdir(exist_ok=True)


class Canvas(BaseModel):
    width: int = Field(gt=0)
    height: int = Field(gt=0)


class Element(BaseModel):
    type: Literal["button", "text", "image", "input", "card"]
    x: float = Field(ge=0)
    y: float = Field(ge=0)
    w: float = Field(gt=0)
    h: float = Field(gt=0)
    text: str | None = None


class Style(BaseModel):
    mode: Literal["lowfi", "hifi"] = "lowfi"
    theme: Literal["gray", "light", "dark"] = "gray"


class FramePage(BaseModel):
    pageId: str
    canvas: Canvas
    elements: list[Element]
    style: Style = Style()


class FrameDocument(BaseModel):
    projectId: str
    pages: list[FramePage]


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/v1/projects/{project_id}/generate")
def generate(project_id: str, payload: FrameDocument) -> dict:
    if payload.projectId != project_id:
        raise HTTPException(status_code=400, detail="projectId mismatch")
    return {
        "projectId": project_id,
        "status": "generated",
        "pages": len(payload.pages),
    }


def _to_inches(px: float, canvas_px: float, slide_inches: float) -> float:
    return (px / canvas_px) * slide_inches


def build_pptx(doc: FrameDocument, out_path: Path) -> None:
    if Presentation is None:
        raise RuntimeError("python-pptx is not available")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide_w_in = prs.slide_width / 914400
    slide_h_in = prs.slide_height / 914400

    for page in doc.pages:
        slide = prs.slides.add_slide(blank_layout)

        for el in page.elements:
            x = Inches(_to_inches(el.x, page.canvas.width, slide_w_in))
            y = Inches(_to_inches(el.y, page.canvas.height, slide_h_in))
            w = Inches(_to_inches(el.w, page.canvas.width, slide_w_in))
            h = Inches(_to_inches(el.h, page.canvas.height, slide_h_in))

            shape = slide.shapes.add_shape(1, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = RGBColor(120, 120, 120)

            if el.type == "button":
                shape.fill.fore_color.rgb = RGBColor(225, 225, 225)

            if el.text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = el.text
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(30, 30, 30)

    prs.save(out_path)


@app.post("/v1/projects/{project_id}/export/pptx")
def export_pptx(project_id: str, payload: FrameDocument):
    if payload.projectId != project_id:
        raise HTTPException(status_code=400, detail="projectId mismatch")

    export_id = f"{project_id}-{uuid.uuid4().hex[:8]}"
    out_path = EXPORT_DIR / f"{export_id}.pptx"

    try:
        build_pptx(payload, out_path)
    except RuntimeError as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    return {"exportId": export_id, "download": f"/v1/exports/{export_id}"}


@app.get("/v1/exports/{export_id}")
def download_export(export_id: str):
    file_path = EXPORT_DIR / f"{export_id}.pptx"
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="export not found")
    return FileResponse(
        path=file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=file_path.name,
    )
